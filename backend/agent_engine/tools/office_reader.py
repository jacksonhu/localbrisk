"""
Office365 file reading tool.

Supported formats:
- Excel (.xlsx, .xls) — via openpyxl
- Word (.docx) — via MarkItDown / docling / python-docx
- PowerPoint (.pptx) — via MarkItDown / docling / python-pptx
- PDF (.pdf) — via MarkItDown / docling

Strategy (three-level fallback):
1. Excel prefers openpyxl (better structured-data fidelity)
2. Word / PowerPoint / PDF:
   a. Fast mode (default): use MarkItDown for quick text summary extraction
   b. Detailed mode: use docling for deep parsing (preserves complex tables, formulas, and layout)
   c. Fallback: if docling is unavailable, downgrade to python-docx / python-pptx
"""

import logging
import os
from pathlib import Path
from typing import Any, Optional, Type

from langchain_core.tools import BaseTool
from pydantic import BaseModel, Field

logger = logging.getLogger(__name__)

# ──────────────────────── Optional dependency detection ────────────────────────

_OPENPYXL_AVAILABLE = False
_MARKITDOWN_AVAILABLE = False
_DOCLING_AVAILABLE = False
_PYTHON_DOCX_AVAILABLE = False
_PYTHON_PPTX_AVAILABLE = False

try:
    import openpyxl
    _OPENPYXL_AVAILABLE = True
except ImportError:
    pass

try:
    from markitdown import MarkItDown
    _MARKITDOWN_AVAILABLE = True
except ImportError:
    pass

try:
    from docling.document_converter import DocumentConverter
    _DOCLING_AVAILABLE = True
except ImportError:
    pass

try:
    from docx import Document as DocxDocument
    _PYTHON_DOCX_AVAILABLE = True
except ImportError:
    pass

try:
    from pptx import Presentation
    _PYTHON_PPTX_AVAILABLE = True
except ImportError:
    pass

# ──────────────────────── Supported extensions ────────────────────────

EXCEL_EXTENSIONS = {".xlsx", ".xls", ".xlsm", ".xlsb"}
WORD_EXTENSIONS = {".docx", ".doc"}
PPTX_EXTENSIONS = {".pptx", ".ppt"}
PDF_EXTENSIONS = {".pdf"}

ALL_SUPPORTED = EXCEL_EXTENSIONS | WORD_EXTENSIONS | PPTX_EXTENSIONS | PDF_EXTENSIONS


# ══════════════════════════════════════════════════════════════
# Low-level reader functions
# ══════════════════════════════════════════════════════════════

def read_excel(file_path: str, sheet_name: Optional[str] = None,
               max_rows: int = 500) -> str:
    """Read an Excel file with openpyxl and return Markdown table output."""
    if not _OPENPYXL_AVAILABLE:
        raise ImportError("openpyxl is not installed. Run: pip install openpyxl")

    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    target_sheets = [sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.sheetnames

    parts: list[str] = []
    for sname in target_sheets:
        ws = wb[sname]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            parts.append(f"## Sheet: {sname}\n\n(Empty sheet)")
            continue

        # Read header row
        header = rows[0]
        col_names = [str(c) if c is not None else "" for c in header]
        data_rows = rows[1:max_rows + 1]

        # Build Markdown table
        md_lines = [f"## Sheet: {sname}", ""]
        md_lines.append("| " + " | ".join(col_names) + " |")
        md_lines.append("| " + " | ".join(["---"] * len(col_names)) + " |")
        for row in data_rows:
            cells = [str(c) if c is not None else "" for c in row]
            md_lines.append("| " + " | ".join(cells) + " |")

        if len(rows) - 1 > max_rows:
            md_lines.append(f"\n> ⚠️ Showing only the first {max_rows} rows out of {len(rows) - 1} total data rows")

        parts.append("\n".join(md_lines))

    wb.close()

    summary = f"**File**: `{os.path.basename(file_path)}`\n"
    summary += f"**Sheet count**: {len(target_sheets)}\n\n"
    return summary + "\n\n---\n\n".join(parts)


def read_with_markitdown(file_path: str) -> str:
    """Quickly parse documents (Word / PPT / PDF / Excel) with MarkItDown and return Markdown.

    MarkItDown is a lightweight open-source converter from Microsoft.
    It is fast and dependency-light, suitable for quick text summary extraction.
    Fidelity for complex tables and formulas is lower than docling.
    """
    if not _MARKITDOWN_AVAILABLE:
        raise ImportError("markitdown is not installed. Run: pip install markitdown")

    md = MarkItDown()
    result = md.convert(file_path)
    content = result.text_content or ""

    summary = f"**File**: `{os.path.basename(file_path)}`\n"
    summary += f"**Format**: {Path(file_path).suffix.upper()}\n"
    summary += f"**Parser**: MarkItDown (fast mode)\n\n"
    return summary + content


def read_with_docling(file_path: str) -> str:
    """Parse documents (Word / PPT / PDF) with docling and return Markdown."""
    if not _DOCLING_AVAILABLE:
        raise ImportError("docling is not installed. Run: pip install docling")

    converter = DocumentConverter()
    result = converter.convert(file_path)
    md_content = result.document.export_to_markdown()

    summary = f"**File**: `{os.path.basename(file_path)}`\n"
    summary += f"**Format**: {Path(file_path).suffix.upper()}\n"
    summary += f"**Parser**: Docling (detailed mode)\n\n"
    return summary + md_content


def read_docx_fallback(file_path: str) -> str:
    """Read Word documents via python-docx (fallback when docling is unavailable)."""
    if not _PYTHON_DOCX_AVAILABLE:
        raise ImportError("python-docx is not installed. Run: pip install python-docx")

    doc = DocxDocument(file_path)
    parts: list[str] = []

    parts.append(f"**File**: `{os.path.basename(file_path)}`\n")

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        # Map heading levels
        if para.style and para.style.name.startswith("Heading"):
            try:
                level = int(para.style.name.replace("Heading ", "").replace("Heading", "1"))
            except ValueError:
                level = 1
            parts.append(f"{'#' * level} {text}")
        else:
            parts.append(text)

    # Extract tables
    for i, table in enumerate(doc.tables):
        parts.append(f"\n### Table {i + 1}\n")
        for j, row in enumerate(table.rows):
            cells = [cell.text.strip() for cell in row.cells]
            parts.append("| " + " | ".join(cells) + " |")
            if j == 0:
                parts.append("| " + " | ".join(["---"] * len(cells)) + " |")

    return "\n\n".join(parts)


def read_pptx_fallback(file_path: str) -> str:
    """Read PPT files via python-pptx (fallback when docling is unavailable)."""
    if not _PYTHON_PPTX_AVAILABLE:
        raise ImportError("python-pptx is not installed. Run: pip install python-pptx")

    prs = Presentation(file_path)
    parts: list[str] = []
    parts.append(f"**File**: `{os.path.basename(file_path)}`\n")
    parts.append(f"**Slide count**: {len(prs.slides)}\n")

    for idx, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_texts.append("| " + " | ".join(cells) + " |")

        parts.append(f"## Slide {idx}")
        if slide_texts:
            parts.append("\n".join(slide_texts))
        else:
            parts.append("(Blank slide)")

    return "\n\n".join(parts)


# ══════════════════════════════════════════════════════════════
# LangChain tool definitions
# ══════════════════════════════════════════════════════════════

class OfficeReaderInput(BaseModel):
    """Input parameters for the Office file reader tool."""
    file_path: str = Field(
        description="Path to the Office file to read (supports .xlsx/.xls/.docx/.pptx/.pdf)"
    )
    sheet_name: Optional[str] = Field(
        default=None,
        description="(Excel only) Sheet name to read; if omitted, all sheets are read"
    )
    max_rows: int = Field(
        default=500,
        description="(Excel only) Maximum number of data rows to read; default is 500"
    )
    mode: str = Field(
        default="fast",
        description=(
            "Parsing mode: "
            "'fast' = use MarkItDown for quick text summary extraction (default, faster); "
            "'detailed' = use Docling for deep parsing (preserves complex tables, formulas, and layout, slower)"
        )
    )


class OfficeReaderTool(BaseTool):
    """Read Office365 file content (Excel / Word / PowerPoint / PDF).

    Supported formats:
    - Excel (.xlsx, .xls): returns data in Markdown table format
    - Word (.docx): returns document content in Markdown
    - PowerPoint (.pptx): returns slide content in Markdown
    - PDF (.pdf): returns document content in Markdown

    Supported parsing modes:
    - fast (default): use MarkItDown for quick text summary extraction
    - detailed: use Docling for deep parsing while preserving complex tables, formulas, and layout
    """

    name: str = "office_reader"
    description: str = (
        "Read Office file content. Supports Excel (.xlsx/.xls), Word (.docx), "
        "PowerPoint (.pptx), and PDF (.pdf). Returns Markdown-formatted file content. "
        "Supports mode parameter: 'fast' (default) uses MarkItDown for quick text summary extraction; "
        "'detailed' uses Docling for deep parsing of complex tables and formulas. "
        "File path can be a virtual path or a real OS path."
    )
    args_schema: Type[BaseModel] = OfficeReaderInput

    # Runtime-injected CompositeBackend instance for virtual path resolving
    _backend: Any = None

    def _run(
        self,
        file_path: str,
        sheet_name: Optional[str] = None,
        max_rows: int = 500,
        mode: str = "fast",
    ) -> str:
        """Synchronous execution."""
        return self._read_file(file_path, sheet_name, max_rows, mode)

    async def _arun(
        self,
        file_path: str,
        sheet_name: Optional[str] = None,
        max_rows: int = 500,
        mode: str = "fast",
    ) -> str:
        """Asynchronous execution (runs sync I/O in a thread pool)."""
        import asyncio
        return await asyncio.to_thread(self._read_file, file_path, sheet_name, max_rows, mode)

    def _read_file(
        self,
        file_path: str,
        sheet_name: Optional[str] = None,
        max_rows: int = 500,
        mode: str = "fast",
    ) -> str:
        """Core read logic.

        mode:
          - "fast": prefer MarkItDown (quick text summary extraction)
          - "detailed": prefer Docling (deep parsing, preserves table/formula/layout)
        """
        # If backend is available, try resolving virtual path to real OS path
        resolved_path = file_path
        if self._backend is not None:
            try:
                from agent_engine.utils.path_resolver import resolve_virtual_path
                resolved_path = resolve_virtual_path(self._backend, file_path)
                if resolved_path != file_path:
                    logger.info(f"Virtual path resolved: {file_path} -> {resolved_path}")
            except Exception as e:
                logger.warning(f"Virtual path resolution failed, fallback to original path: {file_path}, error={e}")
                resolved_path = file_path

        path = Path(resolved_path)
        # Verify file exists
        if not path.exists():
            return f"❌ File not found: {file_path}" + (
                f"\n(Resolved path: {resolved_path})" if resolved_path != file_path else ""
            )

        if not path.is_file():
            return f"❌ Path is not a file: {file_path}"

        ext = path.suffix.lower()

        if ext not in ALL_SUPPORTED:
            return (
                f"❌ Unsupported file type: {ext}\n"
                f"Supported types: {', '.join(sorted(ALL_SUPPORTED))}"
            )

        use_detailed = (mode == "detailed")

        try:
            # ─── Excel: always use openpyxl (best structured-data fidelity) ───
            if ext in EXCEL_EXTENSIONS:
                if _OPENPYXL_AVAILABLE:
                    return read_excel(resolved_path, sheet_name, max_rows)
                # Excel fallback: MarkItDown can also handle xlsx
                if _MARKITDOWN_AVAILABLE:
                    return read_with_markitdown(resolved_path)
                return self._missing_deps_message("Excel", ["openpyxl", "markitdown"])

            # ─── Word / PowerPoint / PDF: three-level fallback strategy ───
            return self._read_document(resolved_path, ext, use_detailed)

        except Exception as e:
            logger.error(f"Failed to read file: {file_path}, error={e}", exc_info=True)
            return f"❌ Failed to read file: {file_path}\nError: {str(e)}"

    def _read_document(self, resolved_path: str, ext: str, use_detailed: bool) -> str:
        """Read document-type files (Word / PPT / PDF) with a three-level fallback strategy.

        Detailed mode: Docling → MarkItDown → python-docx/pptx
        Fast mode: MarkItDown → Docling → python-docx/pptx
        """
        if use_detailed:
            # Detailed mode: prefer Docling
            if _DOCLING_AVAILABLE:
                return read_with_docling(resolved_path)
            if _MARKITDOWN_AVAILABLE:
                logger.info(f"Docling unavailable, downgrade to MarkItDown: {resolved_path}")
                return read_with_markitdown(resolved_path)
        else:
            # Fast mode: prefer MarkItDown
            if _MARKITDOWN_AVAILABLE:
                return read_with_markitdown(resolved_path)
            if _DOCLING_AVAILABLE:
                logger.info(f"MarkItDown unavailable, downgrade to Docling: {resolved_path}")
                return read_with_docling(resolved_path)

        # Final fallback: python-docx / python-pptx
        if ext in WORD_EXTENSIONS and _PYTHON_DOCX_AVAILABLE and ext == ".docx":
            return read_docx_fallback(resolved_path)
        if ext in PPTX_EXTENSIONS and _PYTHON_PPTX_AVAILABLE and ext == ".pptx":
            return read_pptx_fallback(resolved_path)

        # No available parser
        deps = ["markitdown", "docling"]
        if ext in WORD_EXTENSIONS:
            deps.append("python-docx")
        elif ext in PPTX_EXTENSIONS:
            deps.append("python-pptx")
        file_type = {".pdf": "PDF"}.get(ext, ext.upper().lstrip("."))
        if ext in WORD_EXTENSIONS:
            file_type = "Word"
        elif ext in PPTX_EXTENSIONS:
            file_type = "PowerPoint"
        return self._missing_deps_message(file_type, deps)

    @staticmethod
    def _missing_deps_message(file_type: str, packages: list[str]) -> str:
        cmds = " ".join(f"pip install {p}" for p in packages)
        return (
            f"❌ Unable to read {file_type} file because required dependencies are missing.\n"
            f"Please install: {cmds}"
        )


def create_office_reader_tool() -> OfficeReaderTool:
    """Factory function: create an OfficeReaderTool instance."""
    return OfficeReaderTool()


def get_available_formats() -> dict[str, bool]:
    """Return availability status for each supported format in the current environment."""
    return {
        "Excel (.xlsx/.xls)": _OPENPYXL_AVAILABLE,
        "Word/PPT/PDF [markitdown]": _MARKITDOWN_AVAILABLE,
        "Word (.docx) [docling]": _DOCLING_AVAILABLE,
        "Word (.docx) [python-docx]": _PYTHON_DOCX_AVAILABLE,
        "PowerPoint (.pptx) [docling]": _DOCLING_AVAILABLE,
        "PowerPoint (.pptx) [python-pptx]": _PYTHON_PPTX_AVAILABLE,
        "PDF (.pdf) [docling]": _DOCLING_AVAILABLE,
    }
