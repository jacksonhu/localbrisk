"""Unified local file read/write tools bound to one agent directory."""

from __future__ import annotations

import asyncio
import logging
import mimetypes
import os
from pathlib import Path
from typing import Optional, Type

from langchain_core.tools import BaseTool
from pydantic import BaseModel, Field

logger = logging.getLogger(__name__)

# ──────────────────────── Optional dependency detection ────────────────────────

_OPENPYXL_AVAILABLE = False
_MARKITDOWN_AVAILABLE = False
_DOCLING_AVAILABLE = False
_PYTHON_DOCX_AVAILABLE = False
_PYTHON_PPTX_AVAILABLE = False

try:
    import openpyxl

    _OPENPYXL_AVAILABLE = True
except ImportError:
    pass

try:
    from markitdown import MarkItDown

    _MARKITDOWN_AVAILABLE = True
except ImportError:
    pass

try:
    from docling.document_converter import DocumentConverter

    _DOCLING_AVAILABLE = True
except ImportError:
    pass

try:
    from docx import Document as DocxDocument

    _PYTHON_DOCX_AVAILABLE = True
except ImportError:
    pass

try:
    from pptx import Presentation

    _PYTHON_PPTX_AVAILABLE = True
except ImportError:
    pass

# ──────────────────────── Supported extensions ────────────────────────

EXCEL_EXTENSIONS = {".xlsx", ".xls", ".xlsm", ".xlsb"}
WORD_EXTENSIONS = {".docx", ".doc"}
PPTX_EXTENSIONS = {".pptx", ".ppt"}
PDF_EXTENSIONS = {".pdf"}
OFFICE_EXTENSIONS = EXCEL_EXTENSIONS | WORD_EXTENSIONS | PPTX_EXTENSIONS | PDF_EXTENSIONS

TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".rst",
    ".json",
    ".jsonl",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".cfg",
    ".conf",
    ".properties",
    ".env",
    ".xml",
    ".html",
    ".htm",
    ".xhtml",
    ".svg",
    ".csv",
    ".tsv",
    ".log",
    ".sql",
    ".py",
    ".pyi",
    ".js",
    ".jsx",
    ".ts",
    ".tsx",
    ".java",
    ".kt",
    ".kts",
    ".go",
    ".rs",
    ".c",
    ".cc",
    ".cpp",
    ".cxx",
    ".h",
    ".hpp",
    ".cs",
    ".php",
    ".rb",
    ".swift",
    ".scala",
    ".sh",
    ".bash",
    ".zsh",
    ".ps1",
    ".bat",
    ".vue",
    ".css",
    ".scss",
    ".less",
    ".sass",
    ".dockerfile",
    ".gitignore",
    ".gitattributes",
}

DEFAULT_TEXT_ENCODINGS = ("utf-8", "utf-8-sig", "gb18030")
TEXT_MIME_TYPES = {
    "application/json",
    "application/xml",
    "application/x-yaml",
    "application/yaml",
    "image/svg+xml",
}


# ══════════════════════════════════════════════════════════════
# Low-level file helpers
# ══════════════════════════════════════════════════════════════


def read_excel(file_path: str, sheet_name: Optional[str] = None, max_rows: int = 500) -> str:
    """Read an Excel file with openpyxl and return Markdown table output."""
    if not _OPENPYXL_AVAILABLE:
        raise ImportError("openpyxl is not installed. Run: pip install openpyxl")

    workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    target_sheets = [sheet_name] if sheet_name and sheet_name in workbook.sheetnames else workbook.sheetnames

    parts: list[str] = []
    for current_sheet in target_sheets:
        worksheet = workbook[current_sheet]
        rows = list(worksheet.iter_rows(values_only=True))
        if not rows:
            parts.append(f"## Sheet: {current_sheet}\n\n(Empty sheet)")
            continue

        header = rows[0]
        column_names = [str(cell) if cell is not None else "" for cell in header]
        data_rows = rows[1 : max_rows + 1]

        lines = [f"## Sheet: {current_sheet}", ""]
        lines.append("| " + " | ".join(column_names) + " |")
        lines.append("| " + " | ".join(["---"] * len(column_names)) + " |")
        for row in data_rows:
            cells = [str(cell) if cell is not None else "" for cell in row]
            lines.append("| " + " | ".join(cells) + " |")

        if len(rows) - 1 > max_rows:
            lines.append("")
            lines.append(
                f"> Showing only the first {max_rows} rows out of {len(rows) - 1} total data rows"
            )

        parts.append("\n".join(lines))

    workbook.close()
    return "\n\n---\n\n".join(parts)


def read_with_markitdown(file_path: str) -> str:
    """Quickly parse documents with MarkItDown and return Markdown output."""
    if not _MARKITDOWN_AVAILABLE:
        raise ImportError("markitdown is not installed. Run: pip install markitdown")

    markitdown = MarkItDown()
    result = markitdown.convert(file_path)
    return (result.text_content or "").strip()


def read_with_docling(file_path: str) -> str:
    """Parse documents with docling and return Markdown output."""
    if not _DOCLING_AVAILABLE:
        raise ImportError("docling is not installed. Run: pip install docling")

    converter = DocumentConverter()
    result = converter.convert(file_path)
    return result.document.export_to_markdown().strip()


def read_docx_fallback(file_path: str) -> str:
    """Read Word documents via python-docx."""
    if not _PYTHON_DOCX_AVAILABLE:
        raise ImportError("python-docx is not installed. Run: pip install python-docx")

    document = DocxDocument(file_path)
    parts: list[str] = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        if paragraph.style and paragraph.style.name.startswith("Heading"):
            try:
                level = int(paragraph.style.name.replace("Heading ", "").replace("Heading", "1"))
            except ValueError:
                level = 1
            parts.append(f"{'#' * level} {text}")
        else:
            parts.append(text)

    for index, table in enumerate(document.tables, 1):
        parts.append(f"\n### Table {index}\n")
        for row_index, row in enumerate(table.rows):
            cells = [cell.text.strip() for cell in row.cells]
            parts.append("| " + " | ".join(cells) + " |")
            if row_index == 0:
                parts.append("| " + " | ".join(["---"] * len(cells)) + " |")

    return "\n\n".join(parts).strip()


def read_pptx_fallback(file_path: str) -> str:
    """Read PPT files via python-pptx."""
    if not _PYTHON_PPTX_AVAILABLE:
        raise ImportError("python-pptx is not installed. Run: pip install python-pptx")

    presentation = Presentation(file_path)
    parts: list[str] = [f"Slide count: {len(presentation.slides)}"]

    for index, slide in enumerate(presentation.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_texts.append("| " + " | ".join(cells) + " |")

        parts.append(f"## Slide {index}")
        parts.append("\n".join(slide_texts) if slide_texts else "(Blank slide)")

    return "\n\n".join(parts).strip()


def _normalize_line_endings(text: str) -> str:
    """Normalize mixed line endings to LF for deterministic processing."""
    return text.replace("\r\n", "\n").replace("\r", "\n")


def _detect_line_ending(text: str) -> str:
    """Detect the most appropriate line ending for write-back."""
    if "\r\n" in text:
        return "\r\n"
    return "\n"


def _looks_like_text(raw_bytes: bytes, path: Path) -> bool:
    """Best-effort detection for readable text or code files."""
    if path.suffix.lower() in TEXT_EXTENSIONS:
        return True

    mime_type, _ = mimetypes.guess_type(path.name)
    if mime_type and (mime_type.startswith("text/") or mime_type in TEXT_MIME_TYPES):
        return True

    sample = raw_bytes[:4096]
    if b"\x00" in sample:
        return False

    for encoding in DEFAULT_TEXT_ENCODINGS:
        try:
            sample.decode(encoding)
            return True
        except UnicodeDecodeError:
            continue
    return False


def _decode_text_file(path: Path) -> str:
    """Decode a text file using a small encoding fallback list."""
    raw_bytes = path.read_bytes()
    if not _looks_like_text(raw_bytes, path):
        raise ValueError(f"Unsupported non-text file type for text operations: {path.suffix or '<no extension>'}")

    for encoding in DEFAULT_TEXT_ENCODINGS:
        try:
            return raw_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw_bytes.decode("utf-8", errors="replace")


def _slice_lines(content: str, start_line: int = 1, end_line: Optional[int] = None) -> tuple[str, int, int, int]:
    """Slice content by 1-based inclusive line range."""
    if start_line < 1:
        raise ValueError("start_line must be greater than or equal to 1")
    if end_line is not None and end_line < start_line:
        raise ValueError("end_line must be greater than or equal to start_line")

    normalized = _normalize_line_endings(content)
    lines = normalized.splitlines()
    total_lines = len(lines)

    if total_lines == 0:
        return "", 0, 0, 0
    if start_line > total_lines:
        raise ValueError(f"start_line {start_line} exceeds total lines {total_lines}")

    actual_end = total_lines if end_line is None else min(end_line, total_lines)
    selected = lines[start_line - 1 : actual_end]
    return "\n".join(selected), start_line, actual_end, total_lines


def _render_range_result(
    *,
    path: Path,
    content: str,
    file_kind: str,
    parser_name: str,
    start_line: int,
    end_line: Optional[int],
) -> str:
    """Render a normalized file read result with metadata and optional slicing."""
    selected, actual_start, actual_end, total_lines = _slice_lines(content, start_line=start_line, end_line=end_line)
    summary = [
        f"**File**: `{path.name}`",
        f"**Resolved path**: `{path}`",
        f"**Type**: {file_kind}",
        f"**Parser**: {parser_name}",
        f"**Selected lines**: {actual_start}-{actual_end} / {total_lines}",
        "",
    ]
    if total_lines == 0:
        summary.append("(Empty file)")
    else:
        summary.append(selected)
    return "\n".join(summary)


# ══════════════════════════════════════════════════════════════
# Tool input models
# ══════════════════════════════════════════════════════════════


class FileReadInput(BaseModel):
    """Input parameters for the local file read tool."""

    path: str = Field(
        description=(
            "Path to the file to read. Relative paths are resolved from the configured base directory. "
            "Supports Office files and common text/code files."
        )
    )
    start_line: int = Field(default=1, ge=1, description="1-based start line to read, inclusive")
    end_line: Optional[int] = Field(default=None, ge=1, description="1-based end line to read, inclusive")
    sheet_name: Optional[str] = Field(
        default=None,
        description="Excel only. Sheet name to read. If omitted, all sheets are read.",
    )
    max_rows: int = Field(
        default=500,
        ge=1,
        description="Excel only. Maximum number of data rows to read per sheet.",
    )
    mode: str = Field(
        default="fast",
        description=(
            "Office document parsing mode. 'fast' uses MarkItDown when available. "
            "'detailed' prefers Docling for richer structure extraction."
        ),
    )


class FileWriteInput(BaseModel):
    """Input parameters for the local text file write tool."""

    path: str = Field(
        description=(
            "Path to the text file to modify. Relative paths are resolved from the configured base directory."
        )
    )
    content: str = Field(description="New text content to write or use as the replacement block")
    start_line: Optional[int] = Field(
        default=None,
        ge=1,
        description=(
            "Optional 1-based start line for a replacement range. Omit together with end_line to overwrite "
            "the whole file."
        ),
    )
    end_line: Optional[int] = Field(
        default=None,
        ge=1,
        description=(
            "Optional 1-based end line for a replacement range, inclusive. If omitted while start_line is set, "
            "only one line is replaced."
        ),
    )
    create_if_missing: bool = Field(
        default=True,
        description="Create the target file when it does not exist and no range replacement is requested.",
    )


# ══════════════════════════════════════════════════════════════
# Shared tool base
# ══════════════════════════════════════════════════════════════


class _BaseFileTool(BaseTool):
    """Common path resolution utilities shared by local file tools."""

    _base_path: Optional[str] = None

    def _resolve_file_path(self, path: str) -> Path:
        raw_path = (path or "").strip()
        if not raw_path:
            raise ValueError("path cannot be empty")

        candidate = Path(os.path.expanduser(raw_path))
        if candidate.is_absolute():
            return candidate.resolve()

        if self._base_path:
            base_path = Path(os.path.expanduser(self._base_path)).resolve()
            return (base_path / candidate).resolve()
        return candidate.resolve()

    @staticmethod
    def _missing_deps_message(file_type: str, packages: list[str]) -> str:
        install_commands = " ".join(f"pip install {package}" for package in packages)
        return (
            f"Unable to read {file_type} file because required dependencies are missing.\n"
            f"Please install: {install_commands}"
        )


# ══════════════════════════════════════════════════════════════
# file_read tool
# ══════════════════════════════════════════════════════════════


class FileReadTool(_BaseFileTool):
    """Read office documents, text files, and code files from the local filesystem."""

    name: str = "file_read"
    description: str = (
        "Read local files from the configured base directory. Supports Office files, PDFs, and common text/code "
        "files such as txt, md, yaml, xml, html, json, py, js, ts, java, sql, and more. Supports line-range reading."
    )
    args_schema: Type[BaseModel] = FileReadInput

    def _run(
        self,
        path: str,
        start_line: int = 1,
        end_line: Optional[int] = None,
        sheet_name: Optional[str] = None,
        max_rows: int = 500,
        mode: str = "fast",
    ) -> str:
        return self._read_file(path, start_line, end_line, sheet_name, max_rows, mode)

    async def _arun(
        self,
        path: str,
        start_line: int = 1,
        end_line: Optional[int] = None,
        sheet_name: Optional[str] = None,
        max_rows: int = 500,
        mode: str = "fast",
    ) -> str:
        return await asyncio.to_thread(self._read_file, path, start_line, end_line, sheet_name, max_rows, mode)

    def _read_file(
        self,
        path: str,
        start_line: int,
        end_line: Optional[int],
        sheet_name: Optional[str],
        max_rows: int,
        mode: str,
    ) -> str:
        try:
            resolved_path = self._resolve_file_path(path)
        except ValueError as exc:
            return f"❌ {exc}"

        if not resolved_path.exists():
            return f"❌ File not found: {path}\nResolved path: {resolved_path}"
        if not resolved_path.is_file():
            return f"❌ Path is not a file: {path}"

        try:
            content, file_kind, parser_name = self._load_content(
                path=resolved_path,
                sheet_name=sheet_name,
                max_rows=max_rows,
                mode=mode,
            )
            return _render_range_result(
                path=resolved_path,
                content=content,
                file_kind=file_kind,
                parser_name=parser_name,
                start_line=start_line,
                end_line=end_line,
            )
        except Exception as exc:
            logger.error("Failed to read file %s: %s", path, exc, exc_info=True)
            return f"❌ Failed to read file: {path}\nError: {exc}"

    def _load_content(
        self,
        *,
        path: Path,
        sheet_name: Optional[str],
        max_rows: int,
        mode: str,
    ) -> tuple[str, str, str]:
        extension = path.suffix.lower()
        resolved_path = str(path)
        use_detailed = mode.strip().lower() == "detailed"

        if extension in EXCEL_EXTENSIONS:
            if _OPENPYXL_AVAILABLE:
                return read_excel(resolved_path, sheet_name, max_rows), "excel", "openpyxl"
            if _MARKITDOWN_AVAILABLE:
                return read_with_markitdown(resolved_path), "excel", "markitdown"
            raise RuntimeError(self._missing_deps_message("Excel", ["openpyxl", "markitdown"]))

        if extension in OFFICE_EXTENSIONS:
            return self._read_document(resolved_path, extension, use_detailed)

        return _decode_text_file(path), "text", "text"

    def _read_document(self, resolved_path: str, extension: str, use_detailed: bool) -> tuple[str, str, str]:
        """Read document-type files with a simple fallback strategy."""
        parser_name = ""
        content = ""
        if use_detailed:
            if _DOCLING_AVAILABLE:
                parser_name = "docling"
                content = read_with_docling(resolved_path)
            elif _MARKITDOWN_AVAILABLE:
                logger.info("Docling unavailable, downgrade to MarkItDown: %s", resolved_path)
                parser_name = "markitdown"
                content = read_with_markitdown(resolved_path)
        else:
            if _MARKITDOWN_AVAILABLE:
                parser_name = "markitdown"
                content = read_with_markitdown(resolved_path)
            elif _DOCLING_AVAILABLE:
                logger.info("MarkItDown unavailable, downgrade to Docling: %s", resolved_path)
                parser_name = "docling"
                content = read_with_docling(resolved_path)

        if content:
            return content, extension.lstrip("."), parser_name

        if extension in WORD_EXTENSIONS and _PYTHON_DOCX_AVAILABLE and extension == ".docx":
            return read_docx_fallback(resolved_path), "word", "python-docx"
        if extension in PPTX_EXTENSIONS and _PYTHON_PPTX_AVAILABLE and extension == ".pptx":
            return read_pptx_fallback(resolved_path), "powerpoint", "python-pptx"

        packages = ["markitdown", "docling"]
        file_type = {".pdf": "PDF"}.get(extension, extension.upper().lstrip("."))
        if extension in WORD_EXTENSIONS:
            packages.append("python-docx")
            file_type = "Word"
        elif extension in PPTX_EXTENSIONS:
            packages.append("python-pptx")
            file_type = "PowerPoint"
        raise RuntimeError(self._missing_deps_message(file_type, packages))


# ══════════════════════════════════════════════════════════════
# file_write tool
# ══════════════════════════════════════════════════════════════


class FileWriteTool(_BaseFileTool):
    """Modify local text or code files on disk."""

    name: str = "file_write"
    description: str = (
        "Write or modify local text files in the configured base directory. Supports whole-file overwrite or "
        "line-range replacement for text and code files."
    )
    args_schema: Type[BaseModel] = FileWriteInput

    def _run(
        self,
        path: str,
        content: str,
        start_line: Optional[int] = None,
        end_line: Optional[int] = None,
        create_if_missing: bool = True,
    ) -> str:
        return self._write_file(path, content, start_line, end_line, create_if_missing)

    async def _arun(
        self,
        path: str,
        content: str,
        start_line: Optional[int] = None,
        end_line: Optional[int] = None,
        create_if_missing: bool = True,
    ) -> str:
        return await asyncio.to_thread(self._write_file, path, content, start_line, end_line, create_if_missing)

    def _write_file(
        self,
        path: str,
        content: str,
        start_line: Optional[int],
        end_line: Optional[int],
        create_if_missing: bool,
    ) -> str:
        try:
            resolved_path = self._resolve_file_path(path)
        except ValueError as exc:
            return f"❌ {exc}"

        if end_line is not None and start_line is None:
            return "❌ start_line is required when end_line is provided"
        if start_line is not None and end_line is not None and end_line < start_line:
            return "❌ end_line must be greater than or equal to start_line"

        try:
            if not resolved_path.exists():
                if not create_if_missing:
                    return f"❌ File not found: {path}\nResolved path: {resolved_path}"
                if start_line is not None or end_line is not None:
                    return "❌ Cannot replace a line range in a file that does not exist"
                resolved_path.parent.mkdir(parents=True, exist_ok=True)
                resolved_path.write_text(content, encoding="utf-8")
                logger.info("Created text file via file_write: %s", resolved_path)
                return f"✅ Created file: {resolved_path}"

            if not resolved_path.is_file():
                return f"❌ Path is not a file: {path}"

            original_text = _decode_text_file(resolved_path)
            updated_text = self._apply_text_update(original_text, content, start_line, end_line)
            resolved_path.write_text(updated_text, encoding="utf-8")
            logger.info("Updated text file via file_write: %s", resolved_path)

            if start_line is None and end_line is None:
                return f"✅ Overwrote file: {resolved_path}"

            resolved_end_line = start_line if end_line is None else end_line
            return f"✅ Updated file: {resolved_path}\nReplaced lines: {start_line}-{resolved_end_line}"
        except Exception as exc:
            logger.error("Failed to write file %s: %s", path, exc, exc_info=True)
            return f"❌ Failed to write file: {path}\nError: {exc}"

    def _apply_text_update(
        self,
        original_text: str,
        replacement_text: str,
        start_line: Optional[int],
        end_line: Optional[int],
    ) -> str:
        """Apply overwrite or inclusive line-range replacement to a text payload."""
        if start_line is None and end_line is None:
            return replacement_text

        if start_line is None:
            raise ValueError("start_line is required when replacing a range")

        normalized_original = _normalize_line_endings(original_text)
        normalized_replacement = _normalize_line_endings(replacement_text)
        newline = _detect_line_ending(original_text)
        keep_trailing_newline = original_text.endswith(("\n", "\r"))

        original_lines = normalized_original.splitlines()
        replacement_lines = normalized_replacement.splitlines() if normalized_replacement else []

        total_lines = len(original_lines)
        if start_line > total_lines + 1:
            raise ValueError(f"start_line {start_line} exceeds writable range {total_lines + 1}")

        actual_end_line = start_line if end_line is None else end_line
        if actual_end_line > total_lines and start_line <= total_lines:
            actual_end_line = total_lines

        start_index = start_line - 1
        if start_line == total_lines + 1:
            updated_lines = original_lines + replacement_lines
        else:
            updated_lines = original_lines[:start_index] + replacement_lines + original_lines[actual_end_line:]

        updated_text = newline.join(updated_lines)
        if keep_trailing_newline and updated_text:
            updated_text += newline
        return updated_text


# ══════════════════════════════════════════════════════════════
# Tool factories
# ══════════════════════════════════════════════════════════════


def create_file_read_tool(base_path: Optional[str] = None) -> FileReadTool:
    """Create a file_read tool resolved from one base directory."""
    tool = FileReadTool()
    tool._base_path = base_path
    return tool


def create_file_write_tool(base_path: Optional[str] = None) -> FileWriteTool:
    """Create a file_write tool resolved from one base directory."""
    tool = FileWriteTool()
    tool._base_path = base_path
    return tool


def get_available_formats() -> dict[str, bool]:
    """Return availability status for supported file categories."""
    return {
        "Text / code files": True,
        "Excel (.xlsx/.xls) [openpyxl]": _OPENPYXL_AVAILABLE,
        "Word/PPT/PDF [markitdown]": _MARKITDOWN_AVAILABLE,
        "Word/PPT/PDF [docling]": _DOCLING_AVAILABLE,
        "Word (.docx) [python-docx]": _PYTHON_DOCX_AVAILABLE,
        "PowerPoint (.pptx) [python-pptx]": _PYTHON_PPTX_AVAILABLE,
    }


__all__ = [
    "FileReadInput",
    "FileReadTool",
    "FileWriteInput",
    "FileWriteTool",
    "create_file_read_tool",
    "create_file_write_tool",
    "get_available_formats",
]
